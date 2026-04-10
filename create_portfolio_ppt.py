from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# 색상 팔레트
COLOR_BG_DARK = RGBColor(0x0F, 0x17, 0x2A)       # 진한 네이비
COLOR_ACCENT = RGBColor(0x6C, 0x63, 0xFF)          # 보라
COLOR_ACCENT2 = RGBColor(0x00, 0xD4, 0xAA)         # 민트
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_LIGHT = RGBColor(0xB0, 0xB8, 0xD4)
COLOR_YELLOW = RGBColor(0xFF, 0xD6, 0x00)
COLOR_CARD = RGBColor(0x1A, 0x26, 0x40)

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)
blank_layout = prs.slide_layouts[6]


def add_slide():
    return prs.slides.add_slide(blank_layout)


def bg(slide, color=COLOR_BG_DARK):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def box(slide, l, t, w, h, color, alpha=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def txt(slide, text, l, t, w, h, size=18, bold=False, color=COLOR_WHITE,
        align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txb


def accent_line(slide, l, t, w, color=COLOR_ACCENT):
    box(slide, l, t, w, 0.05, color)


# ─────────────────────────────────────────
# SLIDE 1: 표지
# ─────────────────────────────────────────
s = add_slide(); bg(s)
box(s, 0, 0, 13.33, 7.5, COLOR_BG_DARK)
box(s, 0, 0, 0.4, 7.5, COLOR_ACCENT)
box(s, 0, 5.8, 13.33, 1.7, COLOR_CARD)

txt(s, "itnews-multi-agent", 0.7, 1.2, 12, 1.2, size=52, bold=True, color=COLOR_WHITE)
txt(s, "LangChain + LangGraph 기반 멀티에이전트 블로그 자동화 시스템", 0.7, 2.5, 11, 0.7, size=20, color=COLOR_ACCENT2)
accent_line(s, 0.7, 3.3, 8, COLOR_ACCENT)

txt(s, "📌  GitHub 레포 분석  →  AI 블로그 작성  →  SEO+GEO 최적화  →  자동 발행", 0.7, 3.6, 12, 0.6, size=16, color=COLOR_LIGHT)
txt(s, "GitHub: github.com/meene11/itnews-multi-agent", 0.7, 6.1, 8, 0.4, size=13, color=COLOR_LIGHT)
txt(s, "2026.04.08", 10.5, 6.1, 2.5, 0.4, size=13, color=COLOR_LIGHT, align=PP_ALIGN.RIGHT)

# ─────────────────────────────────────────
# SLIDE 2: 프로젝트 개요
# ─────────────────────────────────────────
s = add_slide(); bg(s)
box(s, 0, 0, 0.4, 7.5, COLOR_ACCENT)
txt(s, "01  프로젝트 개요", 0.7, 0.3, 10, 0.6, size=28, bold=True, color=COLOR_WHITE)
accent_line(s, 0.7, 1.0, 11.5, COLOR_ACCENT)

txt(s, "무엇을 만들었나?", 0.7, 1.2, 6, 0.5, size=18, bold=True, color=COLOR_ACCENT2)
txt(s, "python main.py  명령어 하나로\nIT 뉴스를 자동 수집 → AI 블로그 작성 → dev.to + Hashnode 동시 발행", 0.7, 1.7, 9, 0.9, size=16, color=COLOR_WHITE)

# 카드 3개
cards = [
    ("🎯 멘토 숙제", "멀티에이전트 + 오케스트레이터를\n이용한 자동화 시스템 구현"),
    ("📰 데이터 소스", "it-news-pipeline이 매일 수집하는\nIT 뉴스 기사 (Supabase 저장)"),
    ("🚀 최종 결과", "dev.to + Hashnode에\n한국어 기술 블로그 자동 발행"),
]
for i, (title, desc) in enumerate(cards):
    x = 0.7 + i * 4.1
    box(s, x, 2.9, 3.8, 2.2, COLOR_CARD)
    txt(s, title, x+0.2, 3.0, 3.4, 0.5, size=15, bold=True, color=COLOR_ACCENT2)
    txt(s, desc, x+0.2, 3.55, 3.4, 1.3, size=13, color=COLOR_LIGHT)

txt(s, "비유: 취재→작성→편집→발행 하는 기자 팀 전체를 AI로 대체", 0.7, 5.5, 12, 0.5, size=14, color=COLOR_YELLOW)

# ─────────────────────────────────────────
# SLIDE 3: 핵심 개념
# ─────────────────────────────────────────
s = add_slide(); bg(s)
box(s, 0, 0, 0.4, 7.5, COLOR_ACCENT2)
txt(s, "02  핵심 개념 4가지", 0.7, 0.3, 10, 0.6, size=28, bold=True, color=COLOR_WHITE)
accent_line(s, 0.7, 1.0, 11.5, COLOR_ACCENT2)

concepts = [
    ("🤖 AI 에이전트", "스스로 도구를 선택하고\n목표 달성까지 반복 행동하는 AI"),
    ("👥 멀티에이전트", "여러 전문 AI가\n협력하는 시스템"),
    ("🎼 오케스트레이터", "여러 에이전트의 순서/조건/\n재시도를 총괄하는 감독"),
    ("🔧 하네스", "에이전트가 사용할\n도구들을 등록·관리하는 틀"),
]
for i, (title, desc) in enumerate(concepts):
    x = 0.5 + (i % 2) * 6.3
    y = 1.5 + (i // 2) * 2.3
    box(s, x, y, 5.9, 2.0, COLOR_CARD)
    accent_line(s, x, y, 5.9, COLOR_ACCENT if i % 2 == 0 else COLOR_ACCENT2)
    txt(s, title, x+0.2, y+0.15, 5.5, 0.5, size=16, bold=True, color=COLOR_WHITE)
    txt(s, desc, x+0.2, y+0.7, 5.5, 1.0, size=13, color=COLOR_LIGHT)

txt(s, "LangChain = 에이전트 설계 라이브러리  |  LangGraph = 에이전트 흐름 제어 도구", 0.7, 6.8, 12, 0.4, size=13, color=COLOR_LIGHT, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────
# SLIDE 4: 전체 워크플로우
# ─────────────────────────────────────────
s = add_slide(); bg(s)
box(s, 0, 0, 0.4, 7.5, COLOR_YELLOW)
txt(s, "03  전체 워크플로우", 0.7, 0.3, 10, 0.6, size=28, bold=True, color=COLOR_WHITE)
accent_line(s, 0.7, 1.0, 11.5, COLOR_YELLOW)

steps = [
    (COLOR_ACCENT,  "1️⃣  Research Agent",  "GitHub / Supabase\n데이터 수집"),
    (COLOR_ACCENT2, "2️⃣  Writer Agent",    "한국어 블로그\n초안 작성"),
    (COLOR_YELLOW,  "3️⃣  SEO Agent",       "SEO+GEO\n점수 검사"),
    (COLOR_ACCENT,  "4️⃣  Publisher Agent", "dev.to + Hashnode\n동시 발행"),
]
for i, (color, title, desc) in enumerate(steps):
    x = 0.5 + i * 3.1
    box(s, x, 1.5, 2.8, 2.5, COLOR_CARD)
    box(s, x, 1.5, 2.8, 0.08, color)
    txt(s, title, x+0.15, 1.6, 2.5, 0.5, size=13, bold=True, color=color)
    txt(s, desc, x+0.15, 2.2, 2.5, 1.5, size=13, color=COLOR_WHITE)
    if i < 3:
        txt(s, "→", x+2.85, 2.4, 0.3, 0.5, size=20, bold=True, color=COLOR_LIGHT, align=PP_ALIGN.CENTER)

box(s, 1.5, 4.3, 10, 0.08, COLOR_ACCENT)
txt(s, "ORCHESTRATOR (LangGraph) — 전체 흐름 제어", 1.5, 4.0, 10, 0.4, size=14, bold=True, color=COLOR_ACCENT, align=PP_ALIGN.CENTER)

txt(s, "⚡ SEO 점수 70점 미만 → Writer에게 재작성 요청 (최대 2회)", 1.5, 4.8, 10, 0.4, size=14, color=COLOR_YELLOW, align=PP_ALIGN.CENTER)
txt(s, "✅ 최종 발행 완료 → 터미널에 dev.to / Hashnode URL 출력", 1.5, 5.3, 10, 0.4, size=14, color=COLOR_ACCENT2, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────
# SLIDE 5: 에이전트 상세
# ─────────────────────────────────────────
s = add_slide(); bg(s)
box(s, 0, 0, 0.4, 7.5, COLOR_ACCENT)
txt(s, "04  4개 에이전트 역할", 0.7, 0.3, 10, 0.6, size=28, bold=True, color=COLOR_WHITE)
accent_line(s, 0.7, 1.0, 11.5, COLOR_ACCENT)

agents = [
    ("Research", "입력: GitHub URL\n출력: JSON 리서치 결과\n도구: GitHub API, 웹검색"),
    ("Writer",   "입력: 리서치 JSON\n출력: 마크다운 블로그 초안\nAI 온도: 0.7 (창의성 허용)"),
    ("SEO",      "입력: 블로그 초안\n출력: 점수 + 개선된 본문\n체크: SEO 6항목 + GEO 4항목"),
    ("Publisher","입력: 최적화된 초안\n출력: 발행 URL x2\n플랫폼: dev.to + Hashnode"),
]
for i, (name, desc) in enumerate(agents):
    x = 0.5 + (i % 2) * 6.3
    y = 1.3 + (i // 2) * 2.5
    box(s, x, y, 5.9, 2.2, COLOR_CARD)
    txt(s, f"🔹 {name} Agent", x+0.2, y+0.1, 5.5, 0.5, size=15, bold=True, color=COLOR_ACCENT2)
    txt(s, desc, x+0.2, y+0.65, 5.5, 1.4, size=13, color=COLOR_LIGHT)

# ─────────────────────────────────────────
# SLIDE 6: SEO + GEO
# ─────────────────────────────────────────
s = add_slide(); bg(s)
box(s, 0, 0, 0.4, 7.5, COLOR_ACCENT2)
txt(s, "05  SEO + GEO 최적화", 0.7, 0.3, 10, 0.6, size=28, bold=True, color=COLOR_WHITE)
accent_line(s, 0.7, 1.0, 11.5, COLOR_ACCENT2)

box(s, 0.7, 1.2, 5.7, 4.8, COLOR_CARD)
txt(s, "SEO (검색엔진 최적화)", 0.9, 1.3, 5.3, 0.5, size=16, bold=True, color=COLOR_ACCENT)
seo_items = ["① 제목에 핵심 키워드 포함", "② H2/H3 소제목 구조", "③ 키워드 밀도 2~5%", "④ 본문 1500자 이상", "⑤ 도입부 핵심 요약", "⑥ 태그 4개 이하"]
for i, item in enumerate(seo_items):
    txt(s, item, 0.9, 1.85 + i*0.55, 5.3, 0.45, size=13, color=COLOR_LIGHT)

box(s, 6.9, 1.2, 5.7, 4.8, COLOR_CARD)
txt(s, "GEO (AI 검색 최적화) ← 신개념!", 7.1, 1.3, 5.3, 0.5, size=16, bold=True, color=COLOR_ACCENT2)
geo_items = ["⑦ 첫 문단에 핵심 정보 요약", "⑧ 구체적 수치/사실 포함", "⑨ 소제목을 질문 형태로", "⑩ 코드/예시 명확하게 레이블링"]
for i, item in enumerate(geo_items):
    txt(s, item, 7.1, 1.85 + i*0.55, 5.3, 0.45, size=13, color=COLOR_LIGHT)

txt(s, "GEO = ChatGPT, Perplexity 같은 AI 검색이 내 글을 답변 소스로 활용하도록 최적화", 0.7, 6.3, 12, 0.4, size=13, color=COLOR_YELLOW, align=PP_ALIGN.CENTER)

# ─────────────────────────────────────────
# SLIDE 7: 기술 스택
# ─────────────────────────────────────────
s = add_slide(); bg(s)
box(s, 0, 0, 0.4, 7.5, COLOR_YELLOW)
txt(s, "06  기술 스택", 0.7, 0.3, 10, 0.6, size=28, bold=True, color=COLOR_WHITE)
accent_line(s, 0.7, 1.0, 11.5, COLOR_YELLOW)

stacks = [
    ("AI / LLM", ["Claude Sonnet 4.6 (Anthropic)", "GPT-4o-mini (OpenAI)"]),
    ("에이전트 프레임워크", ["LangChain >= 0.3", "LangGraph >= 0.2"]),
    ("블로그 플랫폼", ["dev.to REST API", "Hashnode GraphQL API"]),
    ("자동화 / 인프라", ["GitHub Actions (스케줄)", "Supabase (뉴스 DB)"]),
    ("검색 도구", ["DuckDuckGo (ddgs)", "GitHub REST API"]),
    ("언어 / 환경", ["Python 3.11+", "pydantic-settings"]),
]
for i, (category, items) in enumerate(stacks):
    x = 0.5 + (i % 3) * 4.2
    y = 1.3 + (i // 3) * 2.4
    box(s, x, y, 3.9, 2.1, COLOR_CARD)
    txt(s, category, x+0.2, y+0.1, 3.5, 0.45, size=13, bold=True, color=COLOR_YELLOW)
    for j, item in enumerate(items):
        txt(s, f"• {item}", x+0.2, y+0.6+j*0.55, 3.5, 0.45, size=12, color=COLOR_LIGHT)

# ─────────────────────────────────────────
# SLIDE 8: Git 브랜치 전략
# ─────────────────────────────────────────
s = add_slide(); bg(s)
box(s, 0, 0, 0.4, 7.5, COLOR_ACCENT)
txt(s, "07  Git 브랜치 전략", 0.7, 0.3, 10, 0.6, size=28, bold=True, color=COLOR_WHITE)
accent_line(s, 0.7, 1.0, 11.5, COLOR_ACCENT)

branches = [
    ("main", "항상 최신 안정 버전 유지", COLOR_ACCENT2),
    ("phase/1-research-agent", "Research Agent 구현 → main merge", COLOR_ACCENT),
    ("phase/2-agents", "Writer / SEO / Publisher 구현 → main merge", COLOR_ACCENT),
    ("phase/3-orchestrator", "LangGraph 오케스트레이터 연결 → main merge", COLOR_ACCENT),
    ("phase/4-automation", "GitHub Actions 자동화 → main merge", COLOR_ACCENT),
    ("feature/geo-optimization", "GEO 항목 추가 → main merge", COLOR_YELLOW),
    ("test/openai-version", "OpenAI 전환 테스트 브랜치 (진행중)", COLOR_ACCENT2),
]
for i, (branch, desc, color) in enumerate(branches):
    y = 1.3 + i * 0.75
    box(s, 0.7, y, 3.5, 0.55, color)
    txt(s, branch, 0.9, y+0.05, 3.3, 0.45, size=11, bold=True, color=COLOR_BG_DARK)
    txt(s, desc, 4.4, y+0.05, 8.5, 0.45, size=12, color=COLOR_LIGHT)

# ─────────────────────────────────────────
# SLIDE 9: 실행 결과
# ─────────────────────────────────────────
s = add_slide(); bg(s)
box(s, 0, 0, 0.4, 7.5, COLOR_ACCENT2)
txt(s, "08  실제 실행 결과", 0.7, 0.3, 10, 0.6, size=28, bold=True, color=COLOR_WHITE)
accent_line(s, 0.7, 1.0, 11.5, COLOR_ACCENT2)

box(s, 0.7, 1.2, 11.5, 3.8, COLOR_CARD)
log = """\
[1/4] Research Agent 실행 중...       ✅ 완료
[2/4] Writer Agent 실행 중...          ✅ 완료
[3/4] SEO Agent 실행 중...
      SEO 점수 85점 — 기준 통과, 발행합니다.  ✅
[4/4] Publisher Agent 실행 중...
      [dev.to] 발행 중...
      [Hashnode] 발행 중...

======================================
  파이프라인 완료!
  제목    : AI 기반 IT 뉴스 데이터 파이프라인 구축 방법
  SEO 점수: 85점
  태그    : ['ai', 'pipeline', 'sentiment', 'trending']"""
txt(s, log, 1.0, 1.3, 11.0, 3.5, size=12, color=COLOR_ACCENT2)

txt(s, "⚠️ 발행 실패 이슈 → 한글 태그 거절 → 영문 자동 변환 로직 추가로 해결", 0.7, 5.3, 11.5, 0.4, size=13, color=COLOR_YELLOW)
txt(s, "⚠️ Anthropic 529 과부하 → 자동 재시도 로직 (10초/20초/30초) 추가로 해결", 0.7, 5.8, 11.5, 0.4, size=13, color=COLOR_YELLOW)

# ─────────────────────────────────────────
# SLIDE 10: 포트폴리오 포인트
# ─────────────────────────────────────────
s = add_slide(); bg(s)
box(s, 0, 0, 0.4, 7.5, COLOR_ACCENT)
txt(s, "09  포트폴리오 핵심 포인트", 0.7, 0.3, 10, 0.6, size=28, bold=True, color=COLOR_WHITE)
accent_line(s, 0.7, 1.0, 11.5, COLOR_ACCENT)

points = [
    ("🏗️ 멀티에이전트 설계", "단일 AI가 아닌 4개의 전문 에이전트로 역할 분리\n→ 확장성·유지보수성 극대화"),
    ("🎼 LangGraph 오케스트레이터", "조건 분기 + 재시도 루프를 State Graph로 구현\n→ 복잡한 워크플로우 자동 제어"),
    ("📈 SEO + GEO 동시 적용", "검색엔진 최적화 + AI 검색 최적화를 동시에\n→ 최신 트렌드 반영"),
    ("🌐 멀티 플랫폼 자동 발행", "dev.to + Hashnode API 연동으로 동시 발행\n→ 실제 서비스에서 동작하는 수준"),
    ("⚙️ GitHub Actions 자동화", "매주 자동 실행 스케줄 설정\n→ 완전 무인 자동화"),
    ("🔁 에러 처리 & 복구", "API 과부하 자동 재시도, 무한 루프 방지\n→ 프로덕션 수준의 안정성"),
]
for i, (title, desc) in enumerate(points):
    x = 0.5 + (i % 2) * 6.3
    y = 1.3 + (i // 2) * 2.0
    box(s, x, y, 5.9, 1.8, COLOR_CARD)
    txt(s, title, x+0.2, y+0.1, 5.5, 0.45, size=14, bold=True, color=COLOR_ACCENT2)
    txt(s, desc, x+0.2, y+0.6, 5.5, 1.0, size=12, color=COLOR_LIGHT)

# ─────────────────────────────────────────
# SLIDE 11: 다음 작업 & 마무리
# ─────────────────────────────────────────
s = add_slide(); bg(s)
box(s, 0, 0, 0.4, 7.5, COLOR_ACCENT2)
txt(s, "10  다음 작업 & 마무리", 0.7, 0.3, 10, 0.6, size=28, bold=True, color=COLOR_WHITE)
accent_line(s, 0.7, 1.0, 11.5, COLOR_ACCENT2)

txt(s, "🔜  내일 할 작업", 0.7, 1.2, 6, 0.5, size=18, bold=True, color=COLOR_YELLOW)
nexts = [
    "Research Agent를 Supabase 뉴스 데이터 조회로 교체",
    "실제 수집된 IT 뉴스 기사로 '이번 주 IT 뉴스 트렌드' 블로그 자동 발행",
    "Anthropic Claude 서버 안정화 후 메인 브랜치도 Claude로 복원",
]
for i, item in enumerate(nexts):
    txt(s, f"  {i+1}.  {item}", 0.7, 1.85+i*0.6, 12, 0.5, size=14, color=COLOR_LIGHT)

txt(s, "📂  프로젝트 링크", 0.7, 3.6, 6, 0.5, size=18, bold=True, color=COLOR_ACCENT2)
links = [
    "itnews-multi-agent:   github.com/meene11/itnews-multi-agent",
    "it-news-pipeline:   github.com/meene11/it-news-pipeline",
    "뉴스 파이프라인 배포: it-news-pipeline.vercel.app",
]
for i, link in enumerate(links):
    txt(s, f"  •  {link}", 0.7, 4.2+i*0.55, 12, 0.45, size=14, color=COLOR_LIGHT)

box(s, 0.7, 6.0, 11.5, 1.1, COLOR_CARD)
txt(s, "비전공자로서 AI 멀티에이전트 + 오케스트레이터 + 자동화 파이프라인을 직접 구현한 프로젝트", 0.9, 6.1, 11.1, 0.4, size=14, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
txt(s, "LangChain · LangGraph · Anthropic · OpenAI · GitHub Actions · Supabase · dev.to · Hashnode", 0.9, 6.55, 11.1, 0.4, size=12, color=COLOR_ACCENT2, align=PP_ALIGN.CENTER)

# 저장
output_path = r"C:\dev\itnews_multi_agent\auto_multi_agent_portfolio.pptx"
prs.save(output_path)
print(f"PPT 저장 완료: {output_path}")
